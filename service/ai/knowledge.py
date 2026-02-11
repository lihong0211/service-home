#!/usr/bin/env python3
"""
知识库模块：多格式资料上传、向量索引与文档管理。

- 知识库 CRUD、文档与分类管理
- 支持 PDF/DOCX/PPTX/TXT/MD 解析；分段策略为固定长度分片（chunk_size/overlap），DOCX 全文合并后分片，不保留段落/表格边界
- 分段预览：仅解析返回列表，不落库；向量化接口（最后一步）再落库并建索引
- 向量库能力委托 service.ai.vector_db
"""

import os
import re
import shutil
import tempfile
from datetime import datetime

from service.ai.files import convert_doc_to_docx_with_libreoffice, convert_ppt_to_pptx_with_libreoffice
from service.ai.vector_db import (
    DB_NAME_PATTERN,
    _append_documents_to_mysql,
    _create_empty_vector_db_on_disk,
    _delete_vector_db_from_disk,
    _rebuild_vector_db_index,
    _save_documents_to_mysql,
    _sync_categories_from_documents,
    append_documents_batch,
    create_vector_db,
    delete_vector_db_by_id,
    get_vector_db_detail,
    load_vector_db,
    list_vector_dbs_from_mysql,
    rebuild_vector_db_from_mysql,
    sync_vector_db_from_disk,
)


def _knowledge_base_storage_root() -> str:
    """知识库上传文件存储根目录：data/knowledge_base/，与 data/vector_dbs 同级。"""
    return os.path.abspath(os.path.join(os.getcwd(), "data", "knowledge_base"))


def _save_upload_to_kb_folder(kb_id: int, tmp_path: str, original_filename: str) -> str:
    """
    将临时文件保存到 data/knowledge_base/{kb_id}/ 下，文件名带时间戳防覆盖。
    返回保存后的绝对路径。
    """
    root = _knowledge_base_storage_root()
    dir_path = os.path.join(root, str(kb_id))
    os.makedirs(dir_path, exist_ok=True)
    safe = re.sub(r"[^\w\u4e00-\u9fff\-\.]", "_", (original_filename or "file").strip()).strip("._") or "file"
    safe = safe[:120]
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    name = f"{stamp}_{safe}"
    dest = os.path.join(dir_path, name)
    shutil.copy2(tmp_path, dest)
    return os.path.abspath(dest)


def vectorize_knowledge_base(knowledge_base_id: int) -> dict:
    """
    将知识库下全部分段向量化并写入关联向量库（无则创建）。
    返回 {"vector_db_id": int, "count": int, "created": bool }。
    """
    from model.ai import (
        KnowledgeBase,
        KnowledgeBaseDocument,
        KnowledgeBaseSegment,
        VectorDb,
    )
    kb = KnowledgeBase.get_by_id(knowledge_base_id)
    if not kb:
        raise FileNotFoundError("知识库不存在")
    docs = KnowledgeBaseDocument.select_by({"knowledge_base_id": knowledge_base_id})
    document_ids = [d.id for d in docs]
    doc_id_to_file_name = {d.id: (d.file_name or "").strip() or "未命名" for d in docs}
    if not document_ids:
        segments = []
    else:
        segments = KnowledgeBaseSegment.select_by({
            "document_id": document_ids,
            "order_by": [{"col": "document_id", "sort": "asc"}, {"col": "index", "sort": "asc"}],
        })
    documents = []
    for seg in segments:
        file_name = doc_id_to_file_name.get(seg.document_id, "未命名")
        meta = {
            "source": "knowledge_base",
            "segment_id": seg.id,
            "document_id": seg.document_id,
            "index": seg.index,
            "parent_id": seg.parent_id,
            "file_name": file_name,
        }
        if getattr(seg, "segment_metadata", None) and isinstance(seg.segment_metadata, dict):
            meta.update(seg.segment_metadata)
        documents.append({
            "id": str(seg.id),
            "text": seg.text or "",
            "category": file_name or f"第{seg.index}段",
            "metadata": meta,
        })
    created = False
    db_name = f"kb_{knowledge_base_id}"
    # 若已有 vector_db_id，优先增量：只对新分段调 embedding，有删除时再全量重建
    if kb.vector_db_id:
        vec_row = VectorDb.get_by_id(kb.vector_db_id)
        if vec_row:
            try:
                db = load_vector_db(vec_row.name)
                existing_ids = {str(d.get("id", "")) for d in (db.get("metadata") or [])}
            except Exception:
                existing_ids = set()
            current_ids = {d["id"] for d in documents}
            new_ids = current_ids - existing_ids
            removed_ids = existing_ids - current_ids
            if removed_ids:
                # 有分段被删，必须全量重建以保持索引一致
                _rebuild_vector_db_index(kb.vector_db_id, documents)
                return {"vector_db_id": kb.vector_db_id, "count": len(documents), "created": False, "incremental": False}
            if new_ids:
                new_docs = [d for d in documents if d["id"] in new_ids]
                appended = append_documents_batch(kb.vector_db_id, new_docs)
                return {"vector_db_id": kb.vector_db_id, "count": appended, "created": False, "incremental": True}
            return {"vector_db_id": kb.vector_db_id, "count": 0, "created": False, "incremental": True}
    create_vector_db(db_name, documents)
    # create_vector_db 只写磁盘；必须把 vector_db 表写入并回填 knowledge_base.vector_db_id（不依赖 sync 从磁盘读，直接用内存 documents）
    row = VectorDb.select_one_by({"name": db_name})
    if row:
        vector_db_id = row.id
        _save_documents_to_mysql(vector_db_id, documents)
        _sync_categories_from_documents(vector_db_id, documents)
    else:
        vector_db_id = VectorDb.insert({"name": db_name, "description": None})
        _save_documents_to_mysql(vector_db_id, documents)
        _sync_categories_from_documents(vector_db_id, documents)
    from app.app import db
    kb_row = KnowledgeBase.get_by_id(knowledge_base_id)
    if kb_row:
        kb_row.vector_db_id = vector_db_id
        db.session.add(kb_row)
        db.session.commit()
    created = True
    return {"vector_db_id": vector_db_id, "count": len(documents), "created": created}


# --------------- PDF 建库 ---------------

def _extract_pdf_pages(pdf_path: str) -> list[tuple[int, str]]:
    """从 PDF 按页提取文本，返回 [(page_num, page_text), ...]，page_num 从 1 开始。"""
    from PyPDF2 import PdfReader
    reader = PdfReader(pdf_path)
    out = []
    for i, page in enumerate(reader.pages, start=1):
        t = page.extract_text()
        out.append((i, (t or "").strip()))
    return out


def _chunk_text(text: str, chunk_size: int = 1000, chunk_overlap: int = 200) -> list[str]:
    """固定长度分片：按 chunk_size 切块，块间重叠 chunk_overlap 字符。默认分段策略。"""
    if not text or chunk_size <= 0:
        return []
    if len(text) <= chunk_size:
        return [text] if text.strip() else []
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end]
        if chunk.strip():
            chunks.append(chunk)
        start = end - chunk_overlap
        if start >= len(text):
            break
    return chunks


def _documents_from_pdf(
    pdf_path: str,
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
) -> list[dict]:
    """
    从 PDF 生成文档列表：按页提取文本，每页内再按 chunk_size/overlap 分块，category 为「第N页」。
    每项 {"id": "p{page}_c{i}", "text": str, "category": "第N页"}。
    """
    pages = _extract_pdf_pages(pdf_path)
    documents = []
    for page_num, page_text in pages:
        if not page_text:
            continue
        chunks = _chunk_text(page_text, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        category = f"第{page_num}页"
        for i, chunk in enumerate(chunks):
            doc_id = f"p{page_num}_c{i}" if len(chunks) > 1 else f"p{page_num}"
            documents.append({"id": doc_id, "text": chunk, "category": category})
    return documents


def _documents_from_docx(
    docx_path: str,
    source_name: str = None,
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
) -> list[dict]:
    """
    从 DOCX 提取全文（段落+表格）后按固定长度分片。
    - category 为来源文件名（source_name）
    - 不保留段落/表格边界信息
    - 返回 [{"id": "{source}_c{i}", "text": str, "category": str}, ...]
    """
    from docx import Document as DocxDocument
    doc = DocxDocument(docx_path)
    source = source_name or "docx"
    parts = []
    for para in doc.paragraphs:
        t = (para.text or "").strip()
        if t:
            parts.append(t)
    for table in doc.tables:
        rows = []
        for row in table.rows:
            rows.append(" | ".join((c.text or "").strip() for c in row.cells))
        t = "\n".join(rows).strip()
        if t:
            parts.append(t)
    full_text = "\n\n".join(parts)
    if not full_text.strip():
        return []
    chunks = _chunk_text(full_text, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    return [
        {"id": f"{source}_c{i}", "text": c, "category": source}
        for i, c in enumerate(chunks)
    ]


def _ocr_image_to_text(image_path: str) -> str | None:
    """
    使用 Tesseract 对图片做 OCR，返回识别文本；失败或未安装返回 None。
    批量上传时已优化：图片缩放 + 快速 PSM 模式。
    """
    try:
        import pytesseract
        from PIL import Image
        
        img = Image.open(image_path)
        # 优化：限制图片尺寸以加速（大图先缩放，批量上传时减少处理时间）
        max_size = 2000
        if max(img.size) > max_size:
            ratio = max_size / max(img.size)
            new_size = (int(img.size[0] * ratio), int(img.size[1] * ratio))
            img = img.resize(new_size, Image.Resampling.LANCZOS)
        
        # 批量上传优化：使用 PSM 6（单列文本块），比默认模式快 30-50%
        # 中英文混合：chi_sim+eng；若未安装 chi_sim 可改为 eng
        try:
            text = pytesseract.image_to_string(img, lang="chi_sim+eng", config="--psm 6")
        except (pytesseract.TesseractError, pytesseract.TesseractNotFoundError):
            try:
                text = pytesseract.image_to_string(img, lang="eng", config="--psm 6")
            except Exception:
                return None
        except Exception:
            return None
        return (text or "").strip() or None
    except Exception:
        return None


def _documents_from_image(
    image_path: str,
    source_name: str = None,
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
) -> list[dict]:
    """从图片 OCR 提取文本后按块切分。无文字时返回空列表。"""
    text = _ocr_image_to_text(image_path)
    if not text:
        return []
    source = source_name or "image"
    chunks = _chunk_text(text, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    return [
        {"id": f"{source}_c{i}", "text": c, "category": source}
        for i, c in enumerate(chunks)
    ]


def _documents_from_txt(
    txt_path: str,
    source_name: str = None,
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
) -> list[dict]:
    """从 TXT/MD 按块切分，category 为文件名。"""
    with open(txt_path, "r", encoding="utf-8", errors="ignore") as f:
        text = f.read()
    source = source_name or "txt"
    chunks = _chunk_text(text, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    return [
        {"id": f"{source}_c{i}", "text": c, "category": source}
        for i, c in enumerate(chunks)
    ]


def _documents_from_pptx(pptx_path: str, source_name: str = None) -> list[dict]:
    """从 PPTX 按页（幻灯片）提取文本，每页内所有形状文本合并为一条，category 为「第N页」。"""
    from pptx import Presentation
    prs = Presentation(pptx_path)
    source = source_name or "pptx"
    documents = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = (para.text or "").strip()
                if t:
                    parts.append(t)
        text = "\n".join(parts)
        if not text:
            continue
        category = f"第{slide_num}页"
        doc_id = f"{source}_s{slide_num}"
        documents.append({"id": doc_id, "text": text, "category": category})
    return documents


def _documents_from_excel(
    excel_path: str,
    source_name: str = None,
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
) -> list[dict]:
    """从 Excel（.xlsx 或 .xls）提取文本：每个工作表按行提取，表格数据转为文本后分块。"""
    import os
    ext = os.path.splitext(excel_path)[1].lower()
    source = source_name or "excel"
    documents = []
    try:
        if ext == ".xlsx":
            try:
                from openpyxl import load_workbook
                wb = load_workbook(excel_path, data_only=True)
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    rows_text = []
                    for row in sheet.iter_rows(values_only=True):
                        row_vals = [str(cell) if cell is not None else "" for cell in row]
                        row_text = "\t".join(row_vals).strip()
                        if row_text:
                            rows_text.append(row_text)
                    if rows_text:
                        full_text = "\n".join(rows_text)
                        category = f"工作表: {sheet_name}"
                        chunks = _chunk_text(full_text, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
                        for i, chunk in enumerate(chunks):
                            doc_id = f"{source}_{sheet_name}_c{i}"
                            documents.append({"id": doc_id, "text": chunk, "category": category})
            except ImportError:
                raise ValueError("解析 .xlsx 需要安装 openpyxl: pip install openpyxl")
        elif ext == ".xls":
            try:
                import xlrd
                wb = xlrd.open_workbook(excel_path)
                for sheet_idx, sheet in enumerate(wb.sheets()):
                    rows_text = []
                    for row_idx in range(sheet.nrows):
                        row_vals = [str(sheet.cell(row_idx, col_idx).value) for col_idx in range(sheet.ncols)]
                        row_text = "\t".join(row_vals).strip()
                        if row_text:
                            rows_text.append(row_text)
                    if rows_text:
                        full_text = "\n".join(rows_text)
                        category = f"工作表: {sheet.name}"
                        chunks = _chunk_text(full_text, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
                        for i, chunk in enumerate(chunks):
                            doc_id = f"{source}_{sheet.name}_c{i}"
                            documents.append({"id": doc_id, "text": chunk, "category": category})
            except ImportError:
                raise ValueError("解析 .xls 需要安装 xlrd: pip install xlrd")
    except Exception as e:
        raise ValueError(f"解析 Excel 文件失败: {e}")
    return documents


def parse_file_to_documents(
    file_path: str,
    filename: str,
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
) -> list[dict]:
    """
    按扩展名解析文件为分段列表（固定长度分片）。
    支持 .pdf / .docx / .pptx / .txt / .md。
    每项 {"id": str, "text": str, "category": str}。
    """
    fn = (filename or "").lower()
    if fn.endswith(".pdf"):
        return _documents_from_pdf(file_path, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    if fn.endswith(".doc"):
        # 用 LibreOffice 将 .doc 转为 .docx 后解析
        tmp_dir = tempfile.mkdtemp(prefix="kb_doc_")
        try:
            docx_path = convert_doc_to_docx_with_libreoffice(file_path, tmp_dir)
            if docx_path and os.path.isfile(docx_path):
                out = _documents_from_docx(
                    docx_path,
                    source_name=filename,
                    chunk_size=chunk_size,
                    chunk_overlap=chunk_overlap,
                )
                return out
        finally:
            try:
                if os.path.isdir(tmp_dir):
                    for f in os.listdir(tmp_dir):
                        try:
                            os.remove(os.path.join(tmp_dir, f))
                        except OSError:
                            pass
                    os.rmdir(tmp_dir)
            except OSError:
                pass
        raise ValueError(
            "旧版 .doc 无法解析：请安装 LibreOffice（如 Mac: brew install --cask libreoffice）后重试，或另存为 .docx 后上传"
        )
    if fn.endswith(".docx"):
        return _documents_from_docx(
            file_path,
            source_name=filename,
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )
    if fn.endswith(".ppt"):
        # 用 LibreOffice 将 .ppt 转为 .pptx 后解析
        tmp_dir = tempfile.mkdtemp(prefix="kb_ppt_")
        try:
            pptx_path = convert_ppt_to_pptx_with_libreoffice(file_path, tmp_dir)
            if pptx_path and os.path.isfile(pptx_path):
                out = _documents_from_pptx(pptx_path, source_name=filename)
                return out
        finally:
            try:
                if os.path.isdir(tmp_dir):
                    for f in os.listdir(tmp_dir):
                        try:
                            os.remove(os.path.join(tmp_dir, f))
                        except OSError:
                            pass
                    os.rmdir(tmp_dir)
            except OSError:
                pass
        raise ValueError(
            "旧版 .ppt 无法解析：请安装 LibreOffice（如 Mac: brew install --cask libreoffice）后重试，或另存为 .pptx 后上传"
        )
    if fn.endswith(".pptx"):
        return _documents_from_pptx(file_path, source_name=filename)
    if fn.endswith(".txt") or fn.endswith(".md"):
        return _documents_from_txt(file_path, source_name=filename, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    if fn.endswith((".xlsx", ".xls")):
        return _documents_from_excel(
            file_path,
            source_name=filename,
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )
    if fn.endswith((".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp")):
        # 免费 OCR（Tesseract）：有则分段，无则仍保存文件、分段数为 0
        # 注意：parse_file_to_documents 无法接收 skip_ocr 参数，这里总是尝试 OCR
        # 如需跳过 OCR，在上传接口层面处理（先保存文件，不调用 parse_file_to_documents）
        docs = _documents_from_image(
            file_path,
            source_name=filename,
            chunk_size=chunk_size,
            chunk_overlap=chunk_overlap,
        )
        if not docs:
            raise ValueError(
                "图片已保存，但未识别到文字（请本机安装 Tesseract：Mac 执行 brew install tesseract tesseract-lang）"
            )
        return docs
    raise ValueError(f"不支持的文件类型: {filename}，仅支持 .pdf / .doc / .docx / .ppt / .pptx / .xls / .xlsx / .txt / .md 或图片")


def _add_document_and_segments_to_kb(
    knowledge_base_id: int,
    file_name: str,
    documents: list[dict],
    file_id: str = None,
    path: str = None,
) -> dict:
    """
    向知识库写入一条文档记录与分段（knowledge_base_document + knowledge_base_segment）。
    返回 {"document_id": int, "segment_count": int}。不操作向量库，需另调 vectorize_knowledge_base。
    """
    from model.ai import KnowledgeBase, KnowledgeBaseDocument, KnowledgeBaseSegment
    kb = KnowledgeBase.get_by_id(knowledge_base_id)
    if not kb:
        raise FileNotFoundError("知识库不存在")
    doc_id = KnowledgeBaseDocument.insert({
        "knowledge_base_id": knowledge_base_id,
        "file_name": file_name or "未命名",
        "path": path,
        "file_id": file_id,
        "status": "segmented",
    })
    for i, doc in enumerate(documents):
        text = (doc.get("text") or doc.get("content") or "").strip()
        if not text:
            continue
        meta = doc.get("metadata") if isinstance(doc.get("metadata"), dict) else None
        seg_row = {"document_id": doc_id, "text": text, "index": i, "parent_id": None}
        if meta is not None:
            seg_row["segment_metadata"] = meta
        KnowledgeBaseSegment.insert(seg_row)
    return {"document_id": doc_id, "segment_count": len([d for d in documents if (d.get("text") or d.get("content") or "").strip()])}


def append_documents_to_kb(db_id: int, documents: list[dict]) -> dict:
    """向已有向量库（vector_db id）追加文档并重建索引。返回当前库 id、name、count、appended。供向量库直连等场景。"""
    from model.ai import VectorDb
    row = VectorDb.get_by_id(db_id)
    if not row:
        raise FileNotFoundError("向量库不存在")
    if not documents:
        return {"id": row.id, "name": row.name, "count": 0, "appended": 0}
    _append_documents_to_mysql(row.id, documents)
    _sync_categories_from_documents(row.id, documents)
    out = rebuild_vector_db_from_mysql(db_id=row.id)
    return {"id": row.id, "name": row.name, "count": out["count"], "appended": len(documents)}


def create_vector_db_from_pdf(
    pdf_path: str,
    name: str,
    description: str = None,
    chunk_size: int = 1000,
    chunk_overlap: int = 200,
) -> dict:
    """
    从 PDF 文件创建向量库：按页分块、生成向量、落盘并写入 MySQL。
    :param pdf_path: PDF 文件路径
    :param name: 向量库名称
    :param description: 库描述
    :param chunk_size: 每块字符数
    :param chunk_overlap: 块间重叠字符数
    :return: 与 create_vector_db 一致，并含 vector_db_id
    """
    from model.ai import VectorDb
    if not DB_NAME_PATTERN.match(name):
        raise ValueError(f"无效的向量库名: {name}")
    if VectorDb.select_one_by({"name": name}):
        raise ValueError(f"库名已存在: {name}")
    documents = _documents_from_pdf(pdf_path, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    if not documents:
        raise ValueError("PDF 中未提取到有效文本")
    row_id = VectorDb.insert({"name": name, "description": description or ""})
    try:
        out = create_vector_db(name, documents)
        docs = out.get("documents") or []
        _save_documents_to_mysql(row_id, docs)
        if docs:
            _sync_categories_from_documents(row_id, docs)
        out["id"] = row_id
        return out
    except Exception as e:
        VectorDb.force_delete({"id": row_id})
        _delete_vector_db_from_disk(name)
        raise e


# --------------- 知识库表读写（按设计 knowledge_base / knowledge_base_document / knowledge_base_segment） ---------------

def _list_knowledge_bases_from_mysql() -> list[dict]:
    """从 MySQL 列出知识库。返回 [{"id", "name", "description", "vector_db_id", "create_at", "update_at", "vector_db_name", "segment_count"}, ...]"""
    from model.ai import KnowledgeBase, KnowledgeBaseDocument, KnowledgeBaseSegment, VectorDb
    rows = KnowledgeBase.select_by({"order_by": [{"col": "id", "sort": "desc"}]})
    items = []
    for r in rows:
        vec_name = None
        if r.vector_db_id:
            v = VectorDb.get_by_id(r.vector_db_id)
            vec_name = v.name if v else None
        doc_ids = [d.id for d in KnowledgeBaseDocument.select_by({"knowledge_base_id": r.id})]
        seg_count = KnowledgeBaseSegment.count({"document_id": doc_ids}) if doc_ids else 0
        items.append({
            "id": r.id,
            "name": r.name,
            "description": (r.description or "").strip() or None,
            "vector_db_id": r.vector_db_id,
            "create_at": r.create_at.isoformat() if r.create_at else None,
            "update_at": r.update_at.isoformat() if r.update_at else None,
            "vector_db_name": vec_name,
            "segment_count": seg_count,
        })
    return items


# --------------- HTTP 接口 ---------------

def list_knowledge_bases_api():
    """列出所有知识库（查 knowledge_base 表）。GET 或 POST 均可。"""
    from flask import jsonify
    items = _list_knowledge_bases_from_mysql()
    return jsonify({"code": 0, "msg": "ok", "data": {"list": items, "names": [x["name"] for x in items]}})


def create_knowledge_base_api():
    """
    新增知识库：仅写 knowledge_base 表（名称、描述、策略等），不建向量库；向量化在「向量化」步骤执行。
    POST body: { "name": "库名", "description": "可选" }
    """
    from flask import request, jsonify
    from model.ai import KnowledgeBase
    data = request.get_json(silent=True) or {}
    name = (data.get("name") or data.get("db") or "").strip()
    if not name:
        raise ValueError("缺少参数 name 或 db（需 POST JSON body）")
    if not DB_NAME_PATTERN.match(name):
        raise ValueError("库名仅允许 a-zA-Z0-9_-")
    description = (data.get("description") or "").strip() or None
    if KnowledgeBase.select_one_by({"name": name}):
        raise ValueError(f"库名已存在: {name}")
    row_id = KnowledgeBase.insert({
        "name": name,
        "description": description,
        "parsing_strategy": data.get("parsing_strategy") or "fast",
        "chunking_strategy": data.get("chunking_strategy") or "custom",
    })
    row = KnowledgeBase.get_by_id(row_id)
    return jsonify({
        "code": 0,
        "msg": "ok",
        "data": {
            "id": row_id,
            "name": name,
            "description": description,
            "vector_db_id": None,
            "segment_count": 0,
        },
    })


def create_knowledge_base_from_pdf_api():
    """
    从 PDF 创建知识库：写 knowledge_base + knowledge_base_document + knowledge_base_segment，并执行向量化。
    POST multipart/form-data: name（必填）, file（PDF 文件）, description/chunk_size/chunk_overlap 可选。
    """
    import tempfile
    from flask import request, jsonify
    from model.ai import KnowledgeBase
    name = (request.form.get("name") or request.form.get("db") or "").strip()
    if not name:
        raise ValueError("缺少参数 name 或 db")
    if not DB_NAME_PATTERN.match(name):
        raise ValueError("库名仅允许 a-zA-Z0-9_-")
    f = request.files.get("file") or request.files.get("pdf")
    if not f or not f.filename:
        raise ValueError("请上传 PDF 文件（字段 file 或 pdf）")
    if not (f.filename or "").lower().endswith(".pdf"):
        raise ValueError("仅支持 PDF 文件")
    description = (request.form.get("description") or "").strip() or None
    if KnowledgeBase.select_one_by({"name": name}):
        raise ValueError(f"库名已存在: {name}")
    try:
        chunk_size = int(request.form.get("chunk_size") or "1000")
        chunk_overlap = int(request.form.get("chunk_overlap") or "200")
    except (TypeError, ValueError):
        chunk_size, chunk_overlap = 1000, 200
    chunk_size = max(100, min(4000, chunk_size))
    chunk_overlap = max(0, min(chunk_size - 1, chunk_overlap))
    tmp_path = None
    try:
        with tempfile.NamedTemporaryFile(suffix=".pdf", delete=False) as tmp:
            tmp_path = tmp.name
            f.save(tmp_path)
        documents = _documents_from_pdf(tmp_path, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        if not documents:
            raise ValueError("PDF 中未提取到有效文本")
        kb_id = KnowledgeBase.insert({"name": name, "description": description, "parsing_strategy": "fast", "chunking_strategy": "custom"})
        fn = (f.filename or "").strip() or "file.pdf"
        saved_path = _save_upload_to_kb_folder(kb_id, tmp_path, fn)
        _add_document_and_segments_to_kb(kb_id, fn, documents, path=saved_path)
        vec = vectorize_knowledge_base(kb_id)
        kb = KnowledgeBase.get_by_id(kb_id)
        return jsonify({
            "code": 0,
            "msg": "ok",
            "data": {
                "id": kb.id,
                "name": name,
                "description": description,
                "count": vec["count"],
            },
        })
    finally:
        if tmp_path and os.path.isfile(tmp_path):
            try:
                os.remove(tmp_path)
            except OSError:
                pass


def preview_segments_from_db_api():
    """
    分段预览（已落库）：按文档 id 列表直接返回库中已有分段，不解析文件。
    GET /ai/knowledge-base/segments/preview?document_ids=1,2,3
    响应：{ "code": 0, "data": { "documents": [ { "document_id", "file_name", "segments": [ { "id", "text", "index", ... } ] }, ... ] } }
    """
    from flask import request, jsonify
    from model.ai import KnowledgeBaseDocument, KnowledgeBaseSegment
    ids_param = request.args.get("document_ids") or request.args.get("document_ids[]")
    if not ids_param:
        raise ValueError("请提供 document_ids（Query，逗号分隔或多次传），如 ?document_ids=1,2,3")
    if isinstance(ids_param, (list, tuple)):
        doc_ids = [int(x) for x in ids_param]
    else:
        doc_ids = [int(x.strip()) for x in str(ids_param).split(",") if x.strip()]
    if not doc_ids:
        raise ValueError("document_ids 不能为空")
    documents_out = []
    for did in doc_ids:
        doc = KnowledgeBaseDocument.get_by_id(did)
        file_name = (doc.file_name or "").strip() if doc else "未命名"
        rows = KnowledgeBaseSegment.select_by({
            "document_id": did,
            "order_by": [{"col": "index", "sort": "asc"}],
        }) if doc else []
        segments = [
            {
                "id": r.id,
                "document_id": r.document_id,
                "text": (r.text or "").strip(),
                "index": getattr(r, "index", 0),
                "parent_id": getattr(r, "parent_id", None),
                "metadata": getattr(r, "segment_metadata", None) if hasattr(r, "segment_metadata") else None,
            }
            for r in rows
        ]
        documents_out.append({"document_id": did, "file_name": file_name, "segments": segments})
    return jsonify({"code": 0, "msg": "ok", "data": {"documents": documents_out}})


def _resegment_one_document(document_id: int, chunk_size: int, chunk_overlap: int) -> int:
    """对单个文档按 path 重新解析并替换分段，返回新分段数量。"""
    from model.ai import KnowledgeBaseDocument, KnowledgeBaseSegment
    doc = KnowledgeBaseDocument.get_by_id(document_id)
    if not doc:
        raise FileNotFoundError(f"文档不存在: {document_id}")
    path = (getattr(doc, "path", None) or "").strip()
    if not path or not os.path.isfile(path):
        raise FileNotFoundError(f"文档文件路径无效或已丢失: document_id={document_id}")
    fn = (doc.file_name or "").strip() or "file"
    documents = parse_file_to_documents(path, fn, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    if not documents:
        return 0
    KnowledgeBaseSegment.force_delete({"document_id": document_id})
    for i, d in enumerate(documents):
        text = (d.get("text") or d.get("content") or "").strip()
        if not text:
            continue
        meta = d.get("metadata") if isinstance(d.get("metadata"), dict) else None
        seg_row = {"document_id": document_id, "text": text, "index": i, "parent_id": None}
        if meta is not None:
            seg_row["segment_metadata"] = meta
        KnowledgeBaseSegment.insert(seg_row)
    return len([d for d in documents if (d.get("text") or d.get("content") or "").strip()])


def execute_segments_api():
    """
    执行分段并落库（不向量化）。使用文档列表 + 分层参数（分段长度、分段重叠）。
    两种用法：
    1）对已入库文档列表批量重新分段：body 传 document_ids（数组）+ 分段长度 chunk_size + 分段重叠 chunk_overlap，对每个文档按 path 重新解析并替换分段。
    2）对已上传文件执行分段并加入知识库：body 传 kb_id + file_id + file_name + 可选 chunk_size/chunk_overlap。
    POST body: JSON，如 { "document_ids": [1, 2], "chunk_size": 1000, "chunk_overlap": 200 }
    """
    from flask import request, jsonify
    from model.ai import KnowledgeBaseDocument, KnowledgeBaseSegment
    data = request.get_json(silent=True) or request.form or {}
    document_ids = data.get("document_ids")
    kb_id = data.get("kb_id")
    file_id = data.get("file_id")
    file_name = (data.get("file_name") or "").strip()
    try:
        chunk_size = int(data.get("chunk_size") or "1000")
        chunk_overlap = int(data.get("chunk_overlap") or "200")
    except (TypeError, ValueError):
        chunk_size, chunk_overlap = 1000, 200
    chunk_size = max(100, min(4000, chunk_size))
    chunk_overlap = max(0, min(chunk_size - 1, chunk_overlap))

    if document_ids is not None:
        # 文档列表 + 分层参数：批量重新分段
        if not isinstance(document_ids, (list, tuple)):
            document_ids = [document_ids] if document_ids != "" else []
        doc_ids = []
        for x in document_ids:
            try:
                doc_ids.append(int(x))
            except (TypeError, ValueError):
                raise ValueError("document_ids 须为数字数组")
        if not doc_ids:
            raise ValueError("document_ids 不能为空")
        results = []
        for did in doc_ids:
            try:
                cnt = _resegment_one_document(did, chunk_size, chunk_overlap)
                results.append({"document_id": did, "segment_count": cnt})
            except FileNotFoundError as e:
                results.append({"document_id": did, "segment_count": 0, "error": str(e)})
        return jsonify({
            "code": 0,
            "msg": "ok",
            "data": {
                "results": results,
                "chunk_size": chunk_size,
                "chunk_overlap": chunk_overlap,
            },
        })

    if kb_id is not None and file_id and file_name:
        # 对已上传文件执行分段并加入知识库（不向量化）
        from service.ai.files import get_file_path
        try:
            kb_id = int(kb_id)
        except (TypeError, ValueError):
            raise ValueError("kb_id 必须为数字")
        file_path = get_file_path(file_id)
        if not file_path or not os.path.isfile(file_path):
            raise FileNotFoundError("文件不存在或已丢失，请先上传文件")
        documents = parse_file_to_documents(file_path, file_name, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        if not documents:
            raise ValueError("文件中未解析出有效内容")
        out = _add_document_and_segments_to_kb(kb_id, file_name, documents, file_id=file_id, path=file_path)
        return jsonify({
            "code": 0,
            "msg": "ok",
            "data": {"document_id": out["document_id"], "segment_count": out["segment_count"], "re_segmented": False},
        })

    raise ValueError("请提供 document_ids（文档列表+分层参数重新分段）或 kb_id+file_id+file_name（对已上传文件执行分段并加入知识库）")


def upload_knowledge_base_api():
    """
    知识库上传资料：仅保存文件到项目目录 data/knowledge_base/{kb_id}/ 并写入
    knowledge_base_document + knowledge_base_segment，不自动向量化；需要检索时再调 POST /ai/knowledge-base/vectorize。
    支持 PDF/DOCX/TXT/MD。可一次传多个 file（多选），或单个 file。可追加到已有知识库（kb_id/kb_name）或新建知识库（name）。
    POST multipart/form-data: file 或 file[]（可多个）, name 或 kb_id/kb_name, description/chunk_size/chunk_overlap 可选。
    """
    import tempfile
    from flask import request, jsonify
    from sqlalchemy.exc import IntegrityError
    from app.app import db
    from model.ai import KnowledgeBase
    # 支持多文件：收集所有名为 file / file[] / file[0],file[1] / files 的表单文件
    file_list = []
    for key in request.files:
        if key in ("file", "file[]", "files") or (key.startswith("file[") and key.endswith("]")):
            file_list.extend(request.files.getlist(key))
    if not file_list:
        single = request.files.get("file")
        if single and single.filename:
            file_list = [single]
    if not file_list or not any((f and (f.filename or "").strip()) for f in file_list):
        raise ValueError("请上传至少一个文件（字段 file / file[] / file[0] / files）")
    _IMAGE_EXT = (".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp")
    allowed = (".pdf", ".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx", ".txt", ".md") + _IMAGE_EXT
    files_to_process = []
    for f in file_list:
        if not f or not f.filename:
            continue
        fn = (f.filename or "").strip()
        if not any(fn.lower().endswith(ext) for ext in allowed):
            continue
        files_to_process.append((f, fn))
    if not files_to_process:
        raise ValueError("请上传至少一个支持的文件（.pdf / .doc / .docx / .ppt / .pptx / .xls / .xlsx / .txt / .md 或图片 .jpg / .png 等）")
    name = (request.form.get("name") or request.form.get("db") or "").strip()
    kb_id = request.form.get("kb_id")
    if kb_id is not None and (kb_id == "" or (isinstance(kb_id, str) and not kb_id.strip())):
        kb_id = None
    kb_name = (request.form.get("kb_name") or request.form.get("kb") or "").strip()
    if not kb_id and not kb_name:
        if not name:
            raise ValueError("新建知识库时请提供 name；或提供 kb_id/kb_name 追加到已有库")
        if not DB_NAME_PATTERN.match(name):
            raise ValueError("库名仅允许 a-zA-Z0-9_-")
        if KnowledgeBase.select_one_by({"name": name}):
            raise ValueError(f"库名已存在: {name}")
        kb_id = None
    else:
        if kb_id is not None:
            try:
                kb_id = int(kb_id)
            except (TypeError, ValueError):
                raise ValueError("kb_id 必须为数字")
            kb = KnowledgeBase.get_by_id(kb_id)
        else:
            kb = KnowledgeBase.select_one_by({"name": kb_name})
            kb_id = kb.id if kb else None
        if not kb:
            raise FileNotFoundError("知识库不存在")
    description = (request.form.get("description") or "").strip() or None
    # 批量上传图片时可选：skip_ocr=true 跳过 OCR（先保存，后续再处理），避免接口超时
    skip_ocr = request.form.get("skip_ocr", "").lower() in ("true", "1", "yes")
    try:
        chunk_size = int(request.form.get("chunk_size") or "1000")
        chunk_overlap = int(request.form.get("chunk_overlap") or "200")
    except (TypeError, ValueError):
        chunk_size, chunk_overlap = 1000, 200
    chunk_size = max(100, min(4000, chunk_size))
    chunk_overlap = max(0, min(chunk_size - 1, chunk_overlap))
    tmp_paths = []
    try:
        if kb_id is None:
            try:
                row_id = KnowledgeBase.insert({"name": name, "description": description, "parsing_strategy": "fast", "chunking_strategy": "custom"})
            except IntegrityError:
                db.session.rollback()
                raise ValueError(f"库名已存在: {name}，请改用 kb_id 追加或换一个库名")
            kb_id = row_id
        kb = KnowledgeBase.get_by_id(kb_id)
        documents_result = []
        # 检查已存在的文档文件名（用于覆盖重复文件）
        from model.ai import KnowledgeBaseDocument, KnowledgeBaseSegment
        existing_docs = KnowledgeBaseDocument.select_by({"knowledge_base_id": kb_id})
        existing_docs_map = {(d.file_name or "").strip().lower(): d for d in existing_docs if d.file_name}
        duplicated_files = []
        for f, fn in files_to_process:
            suffix = ".pdf" if fn.lower().endswith(".pdf") else ".doc" if fn.lower().endswith(".doc") else ".docx" if fn.lower().endswith(".docx") else ".ppt" if fn.lower().endswith(".ppt") else ".pptx" if fn.lower().endswith(".pptx") else ".jpg" if fn.lower().endswith((".jpg", ".jpeg")) else ".png" if fn.lower().endswith(".png") else ".gif" if fn.lower().endswith(".gif") else ".webp" if fn.lower().endswith(".webp") else ".bmp" if fn.lower().endswith(".bmp") else ".txt"
            tmp_path = None
            try:
                # 检查文件名是否重复（忽略大小写），如果重复则删除旧文档
                fn_normalized = (fn or "").strip().lower()
                is_duplicate = False
                if fn_normalized in existing_docs_map:
                    is_duplicate = True
                    duplicated_files.append(fn)
                    old_doc = existing_docs_map[fn_normalized]
                    # 删除旧文档的分段
                    KnowledgeBaseSegment.force_delete({"document_id": old_doc.id})
                    # 删除旧文件（如果存在）
                    old_path = (getattr(old_doc, "path", None) or "").strip()
                    if old_path and os.path.isfile(old_path):
                        try:
                            os.remove(old_path)
                        except OSError:
                            pass
                    # 删除旧文档记录
                    KnowledgeBaseDocument.force_delete({"id": old_doc.id})
                with tempfile.NamedTemporaryFile(suffix=suffix, delete=False) as tmp:
                    tmp_path = tmp.name
                    tmp_paths.append(tmp_path)
                    f.save(tmp_path)
                # 先保存文件到知识库文件夹，即使解析失败也要保存
                saved_path = _save_upload_to_kb_folder(kb_id, tmp_path, fn)
                # 更新已存在列表，避免同批次重复
                if fn_normalized in existing_docs_map:
                    del existing_docs_map[fn_normalized]
                # 批量上传图片时可选：skip_ocr=true 跳过 OCR（先保存，后续再处理）
                is_image = fn.lower().endswith((".jpg", ".jpeg", ".png", ".gif", ".webp", ".bmp"))
                if is_image and skip_ocr:
                    # 跳过 OCR，先保存文件，后续可调用「分段」接口处理
                    out = _add_document_and_segments_to_kb(kb_id, fn, [], path=saved_path)
                    documents_result.append({
                        "file_name": fn,
                        "document_id": out["document_id"],
                        "segment_count": 0,
                        "path": saved_path,
                        "skipped_ocr": True,
                        "overwritten": is_duplicate,
                    })
                    continue
                # 尝试解析文件
                try:
                    documents = parse_file_to_documents(tmp_path, fn, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
                except Exception as parse_err:
                    # 解析失败但文件已保存，创建文档记录（无分段）
                    out = _add_document_and_segments_to_kb(kb_id, fn, [], path=saved_path)
                    documents_result.append({
                        "file_name": fn,
                        "document_id": out["document_id"],
                        "segment_count": 0,
                        "path": saved_path,
                        "error": f"文件解析失败: {str(parse_err)}",
                        "overwritten": is_duplicate,
                    })
                    continue
                if not documents:
                    # 解析返回空（如扫描版PDF），文件已保存，创建文档记录（无分段）
                    out = _add_document_and_segments_to_kb(kb_id, fn, [], path=saved_path)
                    documents_result.append({
                        "file_name": fn,
                        "document_id": out["document_id"],
                        "segment_count": 0,
                        "path": saved_path,
                        "error": "文件中未解析出有效文本（可能是扫描版PDF或图片文件）",
                        "overwritten": is_duplicate,
                    })
                    continue
                # 解析成功，创建文档记录和分段
                out = _add_document_and_segments_to_kb(kb_id, fn, documents, path=saved_path)
                documents_result.append({
                    "file_name": fn,
                    "document_id": out["document_id"],
                    "segment_count": out["segment_count"],
                    "path": saved_path,
                    "overwritten": is_duplicate,
                })
            except Exception as e:
                documents_result.append({"file_name": fn, "document_id": None, "segment_count": 0, "path": None, "error": str(e)})
            finally:
                if tmp_path and os.path.isfile(tmp_path):
                    try:
                        os.remove(tmp_path)
                    except OSError:
                        pass
                if tmp_path in tmp_paths:
                    tmp_paths.remove(tmp_path)
        return jsonify({
            "code": 0,
            "msg": "ok",
            "data": {
                "id": kb.id,
                "name": kb.name,
                "description": getattr(kb, "description", None) or description,
                "documents": documents_result,
                "total": len(documents_result),
                "duplicated_files": duplicated_files if duplicated_files else None,
            },
        })
    except (ValueError, FileNotFoundError):
        raise
    except IntegrityError:
        db.session.rollback()
        raise ValueError("库名已存在或数据冲突，请改用 kb_id 追加或换一个库名")
    except Exception:
        if db.session:
            try:
                db.session.rollback()
            except Exception:
                pass
        raise
    finally:
        for tmp_path in tmp_paths:
            if tmp_path and os.path.isfile(tmp_path):
                try:
                    os.remove(tmp_path)
                except OSError:
                    pass


def list_knowledge_base_documents_api():
    """
    按知识库 id 获取文档列表（用于「分段预览」步骤左侧文档列表）。
    GET /ai/knowledge-base/documents?kb_id=1 或 ?kb_name=xxx
    返回该知识库下所有 knowledge_base_document，含 id、file_name、path、segment_count、status 等。
    """
    from flask import request, jsonify
    from model.ai import KnowledgeBase, KnowledgeBaseDocument, KnowledgeBaseSegment
    kb_id = request.args.get("kb_id")
    kb_name = (request.args.get("kb_name") or request.args.get("kb") or "").strip()
    if not kb_id and not kb_name:
        raise ValueError("请提供 kb_id 或 kb_name（Query 参数）")
    if kb_id is not None:
        try:
            kb_id = int(kb_id)
        except (TypeError, ValueError):
            raise ValueError("kb_id 必须为数字")
        kb = KnowledgeBase.get_by_id(kb_id)
    else:
        kb = KnowledgeBase.select_one_by({"name": kb_name})
        kb_id = kb.id if kb else None
    if not kb:
        raise FileNotFoundError("知识库不存在")
    rows = KnowledgeBaseDocument.select_by({
        "knowledge_base_id": kb_id,
        "order_by": [{"col": "id", "sort": "asc"}],
    })
    doc_ids = [r.id for r in rows]
    seg_counts = {did: KnowledgeBaseSegment.count({"document_id": did}) for did in doc_ids}
    list_ = []
    for r in rows:
        list_.append({
            "id": r.id,
            "knowledge_base_id": r.knowledge_base_id,
            "file_name": (r.file_name or "").strip() or "未命名",
            "path": (r.path or "").strip() or None,
            "file_id": (r.file_id or "").strip() or None,
            "status": (r.status or "").strip() or "pending",
            "segment_count": seg_counts.get(r.id, 0),
            "create_at": r.create_at.isoformat() if getattr(r, "create_at", None) else None,
        })
    return jsonify({"code": 0, "msg": "ok", "data": {"list": list_, "total": len(list_)}})


def get_document_segments_api(document_id: int):
    """
    按文档 id 获取分段列表（用于「分段预览」步骤右侧分段预览面板）。
    GET /ai/knowledge-base/document/<document_id>/segments
    返回该文档下所有 knowledge_base_segment，按 index 排序。
    """
    from flask import jsonify
    from model.ai import KnowledgeBaseDocument, KnowledgeBaseSegment
    doc = KnowledgeBaseDocument.get_by_id(document_id)
    if not doc:
        raise FileNotFoundError("文档不存在")
    rows = KnowledgeBaseSegment.select_by({
        "document_id": document_id,
        "order_by": [{"col": "index", "sort": "asc"}],
    })
    list_ = []
    for r in rows:
        list_.append({
            "id": r.id,
            "document_id": r.document_id,
            "text": (r.text or "").strip(),
            "index": getattr(r, "index", 0),
            "parent_id": getattr(r, "parent_id", None),
            "metadata": getattr(r, "segment_metadata", None) if hasattr(r, "segment_metadata") else None,
        })
    return jsonify({"code": 0, "msg": "ok", "data": {"list": list_, "total": len(list_)}})


def preview_knowledge_document_api(document_id: int):
    """
    GET /ai/knowledge-base/document/<document_id>/preview
    根据知识库文档 id 返回文件预览。文档路径来自 knowledge_base_document.path（上传时保存到 data/knowledge_base/{kb_id}/）。
    PDF 直接返回，DOCX/PPTX 转 PDF 后返回（需 LibreOffice），TXT/MD 返回文本。
    """
    from model.ai import KnowledgeBaseDocument
    from service.ai.files import serve_file_preview
    doc = KnowledgeBaseDocument.get_by_id(document_id)
    if not doc:
        raise FileNotFoundError("文档不存在")
    path = (getattr(doc, "path", None) or "").strip()
    if not path or not os.path.isfile(path):
        raise FileNotFoundError("文件路径无效或文件已丢失")
    cache_dir = os.path.join(_knowledge_base_storage_root(), ".preview_cache")
    return serve_file_preview(
        path,
        (doc.file_name or "").strip() or "file",
        cache_key=f"kb_doc_{document_id}",
        cache_dir=cache_dir,
    )


def sync_knowledge_base_from_disk_api():
    """
    将磁盘上已存在的向量库同步到 MySQL（补写 vector_db、vector_db_document）。
    POST body: { "name": "disney", "description": "可选" }
    """
    from flask import request, jsonify
    data = request.get_json() or {}
    name = (data.get("name") or data.get("db") or "").strip()
    if not name:
        raise ValueError("缺少参数 name 或 db")
    if not DB_NAME_PATTERN.match(name):
        raise ValueError("库名仅允许 a-zA-Z0-9_-")
    description = (data.get("description") or "").strip() or None
    out = sync_vector_db_from_disk(name, description=description)
    return jsonify({"code": 0, "msg": "ok", "data": out})


def rebuild_knowledge_base_api():
    """
    按知识库 id 重建关联的向量库索引（从 MySQL 文档重新生成向量并写回磁盘）。
    POST body: { "id": 1 }（知识库 id）
    """
    from flask import request, jsonify
    from model.ai import KnowledgeBase
    data = request.get_json() or {}
    kb_id = data.get("id")
    if kb_id is None:
        raise ValueError("请提供 id（知识库 id）")
    try:
        kb_id = int(kb_id)
    except (TypeError, ValueError):
        raise ValueError("id 必须为数字")
    kb = KnowledgeBase.get_by_id(kb_id)
    if not kb:
        raise FileNotFoundError("知识库不存在")
    if not kb.vector_db_id:
        raise ValueError("该知识库尚未向量化，请先执行向量化")
    out = rebuild_vector_db_from_mysql(db_id=kb.vector_db_id)
    return jsonify({"code": 0, "msg": "ok", "data": out})


def vectorize_knowledge_base_api():
    """
    按知识库 id 将该库下全部分段向量化并写入关联向量库（无则创建）。
    POST body: { "knowledge_base_id": 1 }
    """
    from flask import request, jsonify
    data = request.get_json() or {}
    kb_id = data.get("knowledge_base_id")
    if kb_id is None:
        raise ValueError("请提供 knowledge_base_id")
    try:
        kb_id = int(kb_id)
    except (TypeError, ValueError):
        raise ValueError("knowledge_base_id 必须为数字")
    out = vectorize_knowledge_base(kb_id)
    return jsonify({"code": 0, "msg": "ok", "data": out})


def vectorize_with_file_api():
    """
    按 file_id 取文件并解析为分段，写入 knowledge_base_document + knowledge_base_segment，再执行向量化。
    请求：POST JSON 或 form — kb_id（知识库 id）, file_id, file_name, chunk_size/chunk_overlap 可选。
    响应：{ "code": 0, "msg": "ok", "data": { "id", "name", "count", "appended" } }
    """
    from flask import request, jsonify
    from model.ai import KnowledgeBase
    from service.ai.files import get_file_path
    data = request.get_json(silent=True) or request.form
    kb_id = data.get("kb_id") or data.get("knowledge_base_id")
    file_id = (data.get("file_id") or "").strip()
    file_name = (data.get("file_name") or "").strip()
    if not kb_id:
        raise ValueError("请提供 kb_id")
    if not file_id:
        raise ValueError("请提供 file_id")
    try:
        kb_id = int(kb_id)
    except (TypeError, ValueError):
        raise ValueError("kb_id 必须为数字")
    file_path = get_file_path(file_id)
    if not file_path or not os.path.isfile(file_path):
        raise FileNotFoundError("文件不存在或已失效")
    if not file_name:
        raise ValueError("请提供 file_name")
    try:
        chunk_size = int(data.get("chunk_size") or "1000")
        chunk_overlap = int(data.get("chunk_overlap") or "200")
    except (TypeError, ValueError):
        chunk_size, chunk_overlap = 1000, 200
    chunk_size = max(100, min(4000, chunk_size))
    chunk_overlap = max(0, min(chunk_size - 1, chunk_overlap))
    documents = parse_file_to_documents(file_path, file_name, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
    if not documents:
        raise ValueError("文件中未解析出有效内容")
    out = _add_document_and_segments_to_kb(kb_id, file_name, documents, file_id=file_id, path=file_path)
    vec = vectorize_knowledge_base(kb_id)
    kb = KnowledgeBase.get_by_id(kb_id)
    return jsonify({
        "code": 0,
        "msg": "ok",
        "data": {"id": kb.id, "name": kb.name, "count": vec["count"], "appended": out["segment_count"]},
    })


def update_knowledge_base_api():
    """
    更新知识库元信息（knowledge_base 表）：name、description、解析/分段策略等；不动文档与向量库。
    POST body: { "id": 1, "name": "可选", "description": "可选", "parsing_strategy", "chunking_strategy", ... }
    """
    from flask import request, jsonify
    from model.ai import KnowledgeBase
    data = request.get_json() or {}
    kb_id = data.get("id")
    if kb_id is None:
        raise ValueError("缺少参数 id")
    try:
        kb_id = int(kb_id)
    except (TypeError, ValueError):
        raise ValueError("id 必须为数字")
    kb = KnowledgeBase.get_by_id(kb_id)
    if not kb:
        raise FileNotFoundError("知识库不存在")
    update_data = {"id": kb_id}
    if data.get("name") is not None:
        name = (data.get("name") or "").strip()
        if not name:
            raise ValueError("name 不能为空")
        if not DB_NAME_PATTERN.match(name):
            raise ValueError("库名仅允许 a-zA-Z0-9_-")
        other = KnowledgeBase.select_one_by({"name": name})
        if other and other.id != kb_id:
            raise ValueError(f"库名已存在: {name}")
        update_data["name"] = name
    if data.get("description") is not None:
        update_data["description"] = (data.get("description") or "").strip() or None
    if data.get("parsing_strategy") is not None:
        update_data["parsing_strategy"] = (data.get("parsing_strategy") or "fast").strip()
    if data.get("chunking_strategy") is not None:
        update_data["chunking_strategy"] = (data.get("chunking_strategy") or "custom").strip()
    if len(update_data) > 1:
        KnowledgeBase.update(update_data)
    row = KnowledgeBase.get_by_id(kb_id)
    return jsonify({
        "code": 0,
        "msg": "ok",
        "data": {"id": row.id, "name": row.name, "description": (row.description or "").strip() or None, "vector_db_id": row.vector_db_id},
    })


def delete_knowledge_base_api():
    """删除知识库：若有关联向量库则先删向量库（MySQL+磁盘），再删 knowledge_base。POST body: { "id": 1 }"""
    from flask import request, jsonify
    from model.ai import KnowledgeBase
    data = request.get_json() or {}
    kb_id = data.get("id")
    if kb_id is None:
        raise ValueError("缺少参数 id")
    try:
        kb_id = int(kb_id)
    except (TypeError, ValueError):
        raise ValueError("id 必须为数字")
    kb = KnowledgeBase.get_by_id(kb_id)
    if not kb:
        raise FileNotFoundError("知识库不存在")
    name = kb.name
    if kb.vector_db_id:
        delete_vector_db_by_id(kb.vector_db_id)
    KnowledgeBase.force_delete({"id": kb_id})
    return jsonify({"code": 0, "msg": "ok", "data": {"id": kb_id, "name": name}})


def delete_knowledge_base_document_api():
    """删除知识库下的一条文档：删除分段、文档记录，并删除磁盘上的文件。POST body: { "document_id": 1 } 或 { "id": 1 }"""
    from flask import request, jsonify
    from model.ai import KnowledgeBaseDocument, KnowledgeBaseSegment
    data = request.get_json() or {}
    document_id = data.get("document_id") or data.get("id")
    if document_id is None:
        raise ValueError("缺少参数 document_id 或 id")
    try:
        document_id = int(document_id)
    except (TypeError, ValueError):
        raise ValueError("document_id 必须为数字")
    doc = KnowledgeBaseDocument.get_by_id(document_id)
    if not doc:
        raise FileNotFoundError(f"文档不存在: {document_id}")
    file_name = (doc.file_name or "").strip() or "未命名"
    path = (getattr(doc, "path", None) or "").strip()
    KnowledgeBaseSegment.force_delete({"document_id": document_id})
    KnowledgeBaseDocument.force_delete({"id": document_id})
    if path and os.path.isfile(path):
        try:
            os.remove(path)
        except OSError:
            pass
    return jsonify({
        "code": 0,
        "msg": "ok",
        "data": {"document_id": document_id, "file_name": file_name},
    })


def get_knowledge_base_detail_api():
    """获取知识库详情（knowledge_base 表）。GET ?id=1 或 ?name=xxx&with_documents=0|1"""
    from flask import request, jsonify
    from model.ai import KnowledgeBase, KnowledgeBaseDocument, KnowledgeBaseSegment, VectorDb
    kb_id = request.args.get("id")
    kb_name = request.args.get("name")
    with_documents = request.args.get("with_documents", "0") in ("1", "true", "yes")
    if not kb_id and not kb_name:
        raise ValueError("请提供 id 或 name")
    if kb_id is not None:
        try:
            kb_id = int(kb_id)
        except (TypeError, ValueError):
            raise ValueError("id 必须为数字")
        kb = KnowledgeBase.get_by_id(kb_id)
    else:
        kb = KnowledgeBase.select_one_by({"name": kb_name})
        kb_id = kb.id if kb else None
    if not kb:
        raise FileNotFoundError("知识库不存在")
    detail = {
        "id": kb.id,
        "name": kb.name,
        "description": (kb.description or "").strip() or None,
        "vector_db_id": kb.vector_db_id,
        "parsing_strategy": (kb.parsing_strategy or "fast").strip(),
        "chunking_strategy": (kb.chunking_strategy or "custom").strip(),
        "create_at": kb.create_at.isoformat() if kb.create_at else None,
        "update_at": kb.update_at.isoformat() if kb.update_at else None,
    }
    if kb.vector_db_id:
        v = VectorDb.get_by_id(kb.vector_db_id)
        detail["vector_db_name"] = v.name if v else None
    else:
        detail["vector_db_name"] = None
    if with_documents:
        docs = KnowledgeBaseDocument.select_by({"knowledge_base_id": kb.id})
        doc_ids = [d.id for d in docs]
        segs = KnowledgeBaseSegment.select_by({"document_id": doc_ids, "order_by": [{"col": "document_id", "sort": "asc"}, {"col": "index", "sort": "asc"}]}) if doc_ids else []
        detail["documents"] = [{"id": d.id, "file_name": d.file_name, "status": d.status} for d in docs]
        detail["segments"] = [{"id": s.id, "document_id": s.document_id, "index": s.index, "text": (s.text or "")[:200]} for s in segs]
    else:
        detail["documents"] = []
        detail["segments"] = []
    return jsonify({"code": 0, "msg": "ok", "data": detail})


if __name__ == "__main__":
    from service.ai import vector_db
    vector_db._ensure_storage()
    print("向量库存储目录:", vector_db._storage_root())
    print("已有向量库:", vector_db.list_vector_dbs())
