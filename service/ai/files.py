#!/usr/bin/env python3
"""
知识库文件存储与预览（本机目录方案）。
- 上传保存到 data/uploads/（可按 kb_id 分子目录）
- 提供列表、预览 URL；DOCX/PPT/PPTX 通过 LibreOffice 转 PDF 后预览
"""

import os
import re
import json
import uuid
import shutil
import subprocess
from flask import request, jsonify, send_file, Response

# 存储根目录，可与 data/vector_dbs 同级
def _upload_root():
    return os.path.abspath(
        os.getenv("UPLOAD_STORAGE", os.path.join(os.getcwd(), "data", "uploads"))
    )

MANIFEST_FILENAME = "manifest.json"
PREVIEW_CACHE_DIR = ".preview_cache"

# 允许上传的扩展名（小写）
ALLOWED_EXT = {".pdf", ".txt", ".md", ".doc", ".docx", ".ppt", ".pptx", ".xls", ".xlsx"}
# 需要转 PDF 才能预览的
CONVERT_TO_PDF_EXT = {".docx", ".doc", ".ppt", ".pptx", ".xls", ".xlsx"}
# 图片扩展名 -> 预览时返回的 mimetype（浏览器内联显示，不触发下载）
IMAGE_PREVIEW_EXT = {
    ".jpg": "image/jpeg",
    ".jpeg": "image/jpeg",
    ".png": "image/png",
    ".gif": "image/gif",
    ".webp": "image/webp",
    ".bmp": "image/bmp",
}


def _manifest_path():
    return os.path.join(_upload_root(), MANIFEST_FILENAME)


def _load_manifest():
    p = _manifest_path()
    if not os.path.isfile(p):
        return {}
    try:
        with open(p, "r", encoding="utf-8") as f:
            return json.load(f)
    except Exception:
        return {}


def _save_manifest(manifest):
    root = _upload_root()
    os.makedirs(root, exist_ok=True)
    p = _manifest_path()
    with open(p, "w", encoding="utf-8") as f:
        json.dump(manifest, f, ensure_ascii=False, indent=2)


def get_file_path(file_id: str) -> str | None:
    """
    根据上传文件返回的 file_id 解析本地存储的绝对路径。
    用于知识库向量化等需按文件落库的场景。
    不存在或 manifest 无 path 时返回 None。
    """
    manifest = _load_manifest()
    meta = manifest.get(file_id)
    if not meta:
        return None
    rel = meta.get("path")
    if not rel:
        return None
    return os.path.join(_upload_root(), rel)


def _safe_filename(name):
    """保留扩展名，仅保留安全字符。"""
    if not name:
        return "file"
    base, ext = os.path.splitext(name)
    safe = re.sub(r"[^\w\u4e00-\u9fff\-\.]", "_", base).strip("._") or "file"
    return safe[:80] + ext.lower()


def _ensure_upload_dirs():
    root = _upload_root()
    os.makedirs(root, exist_ok=True)
    os.makedirs(os.path.join(root, PREVIEW_CACHE_DIR), exist_ok=True)
    return root


def upload_file_api():
    """
    POST /ai/files/upload
    multipart: file（必填）, kb_id（可选，知识库 id，用于分子目录与列表筛选）
    返回：{ code, msg, data: { id, name, size, url, preview_url } }
    不处理上传进度，仅返回完成状态。
    """
    f = request.files.get("file")
    if not f or not f.filename:
        raise ValueError("请上传文件（字段 file）")
    name = (f.filename or "").strip()
    ext = os.path.splitext(name)[1].lower()
    if ext not in ALLOWED_EXT:
        raise ValueError(f"不支持的文件类型，仅支持: {', '.join(sorted(ALLOWED_EXT))}")

    kb_id = request.form.get("kb_id")
    if kb_id is not None and kb_id != "":
        try:
            kb_id = str(int(kb_id))
        except (TypeError, ValueError):
            kb_id = "0"
    else:
        kb_id = "0"

    file_id = uuid.uuid4().hex
    root = _ensure_upload_dirs()
    subdir = os.path.join(root, kb_id)
    os.makedirs(subdir, exist_ok=True)
    safe_name = _safe_filename(name)
    stored_name = f"{file_id}_{safe_name}"
    rel_path = os.path.join(kb_id, stored_name)
    abs_path = os.path.join(root, rel_path)

    try:
        f.save(abs_path)
        size = os.path.getsize(abs_path)
    except Exception as e:
        raise RuntimeError(f"保存文件失败: {e}") from e

    manifest = _load_manifest()
    manifest[file_id] = {
        "path": rel_path,
        "name": name,
        "size": size,
        "kb_id": kb_id,
        "created_at": __import__("datetime").datetime.utcnow().isoformat() + "Z",
    }
    _save_manifest(manifest)

    # 预览 URL：使用相对 API 路径，前端可拼 baseURL
    preview_url = f"/ai/files/{file_id}/preview"
    url = preview_url
    return jsonify({
        "code": 0,
        "msg": "ok",
        "data": {
            "id": file_id,
            "name": name,
            "size": size,
            "url": url,
            "preview_url": preview_url,
        },
    })


def list_files_api():
    """
    GET /ai/files/list?kbId=xxx
    返回：{ code, msg, data: { list: [ { id, name, size, url, preview_url } ] } }
    """
    kb_id = request.args.get("kbId") or request.args.get("kb_id")
    if kb_id is not None and kb_id != "":
        try:
            kb_id = str(int(kb_id))
        except (TypeError, ValueError):
            kb_id = None
    else:
        kb_id = None

    manifest = _load_manifest()
    list_ = []
    for fid, meta in manifest.items():
        if kb_id is not None and meta.get("kb_id") != kb_id:
            continue
        list_.append({
            "id": fid,
            "name": meta.get("name", ""),
            "size": meta.get("size", 0),
            "url": f"/ai/files/{fid}/preview",
            "preview_url": f"/ai/files/{fid}/preview",
        })
    # 按创建时间倒序
    list_.sort(key=lambda x: manifest.get(x["id"], {}).get("created_at", ""), reverse=True)
    return jsonify({"code": 0, "msg": "ok", "data": {"list": list_}})


# macOS: brew install --cask libreoffice 装到 /Applications，命令为 soffice。
# 若出现「无法验证 LibreOffice」弹窗，见 docs/知识库设计/文件存储与预览.md 最后一节。
_LIBREOFFICE_CMD = None


def _get_libreoffice_cmd() -> str | None:
    """返回可用的 LibreOffice 命令行，优先 macOS 应用路径（brew install --cask 装到 /Applications）。"""
    global _LIBREOFFICE_CMD
    if _LIBREOFFICE_CMD is not None:
        return _LIBREOFFICE_CMD
    # macOS 上安装后是 /Applications/LibreOffice.app/Contents/MacOS/soffice，子进程里 --version 可能报 plist 错，只检查存在
    mac_app = "/Applications/LibreOffice.app/Contents/MacOS/soffice"
    if os.path.isfile(mac_app):
        _LIBREOFFICE_CMD = mac_app
        return mac_app
    for cmd in ("libreoffice", "soffice"):
        try:
            subprocess.run([cmd, "--version"], check=True, capture_output=True, timeout=5)
        except (FileNotFoundError, subprocess.CalledProcessError, subprocess.TimeoutExpired):
            continue
        _LIBREOFFICE_CMD = cmd
        return cmd
    return None


def _convert_to_pdf_with_libreoffice(source_path: str, out_dir: str) -> str | None:
    """使用 LibreOffice 将 docx/ppt/pptx/xlsx/xls 转为 PDF。返回生成的 PDF 路径，失败返回 None。"""
    cmd = _get_libreoffice_cmd()
    if not cmd:
        return None
    try:
        # Excel 转 PDF：LibreOffice 默认转换可能会拆分宽表格
        # 可以通过 --infilter 指定过滤器，但 calc_pdf_Export 是默认的，通常不需要显式指定
        # 表格拆分问题主要取决于源文件的页面设置，命令行转换难以完全控制
        subprocess.run(
            [
                cmd,
                "--headless",
                "--convert-to", "pdf",
                "--outdir", out_dir,
                source_path,
            ],
            check=True,
            capture_output=True,
            timeout=60,
        )
        base = os.path.splitext(os.path.basename(source_path))[0]
        pdf_path = os.path.join(out_dir, base + ".pdf")
        return pdf_path if os.path.isfile(pdf_path) else None
    except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
        return None


def convert_ppt_to_pptx_with_libreoffice(ppt_path: str, out_dir: str) -> str | None:
    """使用 LibreOffice 将 .ppt 转为 .pptx，供知识库解析。返回生成的 .pptx 路径，失败返回 None。"""
    cmd = _get_libreoffice_cmd()
    if not cmd:
        return None
    try:
        subprocess.run(
            [
                cmd,
                "--headless",
                "--convert-to", "pptx",
                "--outdir", out_dir,
                ppt_path,
            ],
            check=True,
            capture_output=True,
            timeout=120,
        )
        base = os.path.splitext(os.path.basename(ppt_path))[0]
        pptx_path = os.path.join(out_dir, base + ".pptx")
        return pptx_path if os.path.isfile(pptx_path) else None
    except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
        return None


def convert_doc_to_docx_with_libreoffice(doc_path: str, out_dir: str) -> str | None:
    """使用 LibreOffice 将 .doc 转为 .docx，供知识库解析。返回生成的 .docx 路径，失败返回 None。"""
    cmd = _get_libreoffice_cmd()
    if not cmd:
        return None
    try:
        subprocess.run(
            [
                cmd,
                "--headless",
                "--convert-to", "docx",
                "--outdir", out_dir,
                doc_path,
            ],
            check=True,
            capture_output=True,
            timeout=120,
        )
        base = os.path.splitext(os.path.basename(doc_path))[0]
        docx_path = os.path.join(out_dir, base + ".docx")
        return docx_path if os.path.isfile(docx_path) else None
    except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
        return None


def _docx_to_html_preview(docx_path: str) -> str:
    """使用 python-docx 将 .docx 转为简单 HTML 预览，不依赖 LibreOffice。"""
    from docx import Document as DocxDocument
    import html
    doc = DocxDocument(docx_path)
    parts = []
    for para in doc.paragraphs:
        t = (para.text or "").strip()
        if t:
            parts.append("<p>" + html.escape(t) + "</p>")
    for table in doc.tables:
        parts.append("<div class=\"table-wrap\"><table border=\"1\" cellpadding=\"4\" cellspacing=\"0\">")
        for row in table.rows:
            parts.append("<tr>")
            for c in row.cells:
                parts.append("<td>" + html.escape((c.text or "").strip()) + "</td>")
            parts.append("</tr>")
        parts.append("</table></div>")
    body = "\n".join(parts) if parts else "<p>（无正文）</p>"
    return (
        "<!DOCTYPE html><html><head><meta charset=\"utf-8\"><title>文档预览</title>"
        "<style>body{font-family:system-ui,sans-serif;max-width:800px;margin:1em auto;padding:0 1em;line-height:1.5;} "
        ".table-wrap{overflow-x:auto;margin:0.5em 0;} "
        "table{border-collapse:collapse;min-width:100%;table-layout:auto;} "
        "td{border:1px solid #ddd;padding:4px;}</style></head><body>"
        + body + "</body></html>"
    )


def _pptx_to_html_preview(pptx_path: str) -> str:
    """使用 python-pptx 将 .pptx 转为简单 HTML 预览（按幻灯片），不依赖 LibreOffice。"""
    from pptx import Presentation
    import html
    prs = Presentation(pptx_path)
    parts = []
    for slide_num, slide in enumerate(prs.slides, start=1):
        slide_parts = []
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            for para in shape.text_frame.paragraphs:
                t = (para.text or "").strip()
                if t:
                    slide_parts.append("<p>" + html.escape(t) + "</p>")
        if slide_parts:
            parts.append(f"<section class=\"slide\"><h3>第 {slide_num} 页</h3>\n" + "\n".join(slide_parts) + "</section>")
    body = "\n".join(parts) if parts else "<p>（无正文）</p>"
    return (
        "<!DOCTYPE html><html><head><meta charset=\"utf-8\"><title>演示文稿预览</title>"
        "<style>body{font-family:system-ui,sans-serif;max-width:800px;margin:1em auto;padding:0 1em;line-height:1.5;} .slide{margin:1.5em 0;padding:0.5em 0;border-bottom:1px solid #eee;} .slide h3{font-size:1em;color:#666;}</style></head><body>"
        + body + "</body></html>"
    )


def _excel_to_html_preview(excel_path: str) -> str:
    """使用 openpyxl（.xlsx）或 xlrd（.xls）将 Excel 转为简单 HTML 表格预览，不依赖 LibreOffice。"""
    import html
    ext = os.path.splitext(excel_path)[1].lower()
    parts = []
    try:
        if ext == ".xlsx":
            try:
                from openpyxl import load_workbook
                wb = load_workbook(excel_path, data_only=True)
                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    parts.append(f"<h3>工作表: {html.escape(sheet_name)}</h3>")
                    parts.append("<div class=\"table-wrap\"><table border=\"1\" cellpadding=\"4\" cellspacing=\"0\">")
                    for row in sheet.iter_rows(values_only=True):
                        parts.append("<tr>")
                        for cell in row:
                            val = str(cell) if cell is not None else ""
                            parts.append("<td>" + html.escape(val) + "</td>")
                        parts.append("</tr>")
                    parts.append("</table></div>")
            except ImportError:
                raise ImportError("需要安装 openpyxl: pip install openpyxl")
        elif ext == ".xls":
            try:
                import xlrd
                wb = xlrd.open_workbook(excel_path)
                for sheet_idx, sheet in enumerate(wb.sheets()):
                    parts.append(f"<h3>工作表: {html.escape(sheet.name)}</h3>")
                    parts.append("<div class=\"table-wrap\"><table border=\"1\" cellpadding=\"4\" cellspacing=\"0\">")
                    for row_idx in range(sheet.nrows):
                        parts.append("<tr>")
                        for col_idx in range(sheet.ncols):
                            cell = sheet.cell(row_idx, col_idx)
                            val = str(cell.value) if cell.value else ""
                            parts.append("<td>" + html.escape(val) + "</td>")
                        parts.append("</tr>")
                    parts.append("</table></div>")
            except ImportError:
                raise ImportError("需要安装 xlrd: pip install xlrd")
    except ImportError as e:
        return (
            "<!DOCTYPE html><html><head><meta charset=\"utf-8\"><title>Excel 预览</title>"
            "<style>body{font-family:system-ui,sans-serif;max-width:800px;margin:1em auto;padding:1em;}</style></head><body>"
            f"<p>无法预览 Excel 文件：{html.escape(str(e))}</p>"
            "<p>请安装相应库后重试，或使用 LibreOffice 转 PDF 预览。</p>"
            "</body></html>"
        )
    body = "\n".join(parts) if parts else "<p>（无数据）</p>"
    return (
        "<!DOCTYPE html><html><head><meta charset=\"utf-8\"><title>Excel 预览</title>"
        "<style>body{font-family:system-ui,sans-serif;max-width:1200px;margin:1em auto;padding:0 1em;line-height:1.5;} "
        ".table-wrap{overflow-x:auto;margin:0.5em 0;} "
        "table{border-collapse:collapse;min-width:100%;table-layout:auto;} "
        "td{border:1px solid #ddd;padding:4px;text-align:left;} "
        "h3{margin-top:1.5em;color:#333;}</style></head><body>"
        + body + "</body></html>"
    )


def serve_file_preview(abs_path: str, file_name: str, cache_key: str = None, cache_dir: str = None):
    """
    根据本地文件路径返回预览响应（PDF 直接返回，DOCX/PPTX 转 PDF，TXT/MD 文本）。
    供 /ai/files/<file_id>/preview 与知识库文档预览共用。
    :param cache_key: DOCX/PPTX 转 PDF 时的缓存文件名（不含后缀），如 file_id 或 kb_doc_1
    :param cache_dir: 转 PDF 缓存目录，默认 _upload_root()/.preview_cache
    """
    if not os.path.isfile(abs_path):
        raise FileNotFoundError("文件不存在或已丢失")
    ext = os.path.splitext(file_name or "")[1].lower()
    root = _upload_root()
    cache_dir = cache_dir or os.path.join(root, PREVIEW_CACHE_DIR)
    if ext in CONVERT_TO_PDF_EXT:
        os.makedirs(cache_dir, exist_ok=True)
        cached_pdf = os.path.join(cache_dir, f"{cache_key or 'preview'}.pdf")
        if not os.path.isfile(cached_pdf):
            pdf_path = _convert_to_pdf_with_libreoffice(abs_path, cache_dir)
            if pdf_path:
                if os.path.realpath(pdf_path) != os.path.realpath(cached_pdf):
                    shutil.move(pdf_path, cached_pdf)
            else:
                # DOCX/PPTX 无 LibreOffice 时用 python-docx / python-pptx 做 HTML 预览，避免 macOS 拦截
                if ext == ".docx":
                    try:
                        html_content = _docx_to_html_preview(abs_path)
                        return Response(
                            html_content,
                            mimetype="text/html; charset=utf-8",
                            headers={"X-Preview-Fallback": "html"},
                        )
                    except Exception:
                        pass
                if ext == ".pptx":
                    try:
                        html_content = _pptx_to_html_preview(abs_path)
                        return Response(
                            html_content,
                            mimetype="text/html; charset=utf-8",
                            headers={"X-Preview-Fallback": "html"},
                        )
                    except Exception:
                        pass
                # Excel 无 LibreOffice 时也改为下载（不提供 HTML 降级，保持与 PDF 一致）
                # .doc / .ppt / .xls / .xlsx 或解析失败：改为下载
                resp = send_file(abs_path, as_attachment=True, download_name=file_name or "file")
                resp.headers["X-Preview-Fallback"] = "download"
                resp.headers["X-Preview-Message"] = "本机未安装 LibreOffice，无法在线预览该格式，已改为下载。安装后可预览: brew install --cask libreoffice"
                return resp
        return send_file(
            cached_pdf,
            mimetype="application/pdf",
            as_attachment=False,
            download_name=os.path.splitext(file_name or "file")[0] + ".pdf",
        )
    if ext == ".pdf":
        return send_file(abs_path, mimetype="application/pdf", as_attachment=False)
    if ext in (".txt", ".md"):
        try:
            with open(abs_path, "r", encoding="utf-8", errors="ignore") as f:
                text = f.read()
        except Exception as e:
            raise RuntimeError("读取文件失败") from e
        return Response(text, mimetype="text/plain; charset=utf-8")
    # 图片：以内联方式返回，便于前端 <img src="..."> 或新标签页直接显示
    if ext in IMAGE_PREVIEW_EXT:
        return send_file(
            abs_path,
            mimetype=IMAGE_PREVIEW_EXT[ext],
            as_attachment=False,
            download_name=file_name or "file",
        )
    return send_file(abs_path, as_attachment=True, download_name=file_name or "file")


def preview_file_api(file_id: str):
    """
    GET /ai/files/<file_id>/preview
    直接返回文件流（Content-Type 正确）；DOCX/PPT/PPTX 转为 PDF 后返回（需本机安装 LibreOffice）。
    """
    manifest = _load_manifest()
    meta = manifest.get(file_id)
    if not meta:
        raise FileNotFoundError("文件不存在")
    rel_path = meta.get("path")
    if not rel_path:
        raise FileNotFoundError("文件路径无效")
    root = _upload_root()
    abs_path = os.path.join(root, rel_path)
    return serve_file_preview(abs_path, meta.get("name", "file"), cache_key=file_id)


