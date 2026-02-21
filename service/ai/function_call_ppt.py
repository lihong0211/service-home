"""
MCP PPT 助手：接入低代码/云平台（活字格、阿里云百炼等）

用户输入「基于本月销售数据做个汇报PPT」，AI 自动：
1. 通过 get_business_data 拉取业务数据（Text2SQL）
2. 通过 generate_pptx 生成图文并茂的演示文稿
3. 返回 file_id 与预览/下载地址

可选：若部署了 ChatPPT-MCP（如 YOOTeam/ChatPPT-MCP Streamable HTTP），
可在 tools 中增加 mcpServers 指向该服务，实现更多 PPT 能力（模板、在线编辑等）。
"""

from __future__ import annotations

import json
import os
import uuid
from http import HTTPStatus
from typing import Any, Dict, List, Optional

import dashscope

dashscope.api_key = os.getenv("DASHSCOPE_API_KEY")

# 业务数据查询工具：自然语言 -> 数据库查询结果
GET_BUSINESS_DATA_TOOL = {
    "type": "function",
    "function": {
        "name": "get_business_data",
        "description": "根据自然语言查询业务或销售等数据。例如：本月销售数据、各区域销售额、产品销量统计、年度汇总等。会连接数据库执行查询并返回结果表格的摘要与数据，供后续生成 PPT 使用。",
        "parameters": {
            "type": "object",
            "properties": {
                "question": {
                    "type": "string",
                    "description": "用自然语言描述要查询的数据，如：本月销售数据、各产品销量排行、按区域统计销售额",
                },
            },
            "required": ["question"],
        },
    },
}

# PPT 生成工具：标题 + 幻灯片列表 -> 生成文件并返回 file_id
GENERATE_PPTX_TOOL = {
    "type": "function",
    "function": {
        "name": "generate_pptx",
        "description": "根据汇报标题和幻灯片内容生成 PPT 文件。先根据已有数据或大纲组织 slides，每页包含标题和要点。生成后返回 file_id 和预览/下载地址。",
        "parameters": {
            "type": "object",
            "properties": {
                "title": {
                    "type": "string",
                    "description": "PPT 主标题，通常显示在封面",
                },
                "slides": {
                    "type": "array",
                    "description": "幻灯片列表，除封面外每页的标题与内容",
                    "items": {
                        "type": "object",
                        "properties": {
                            "title": {"type": "string", "description": "本页标题"},
                            "content": {
                                "type": "string",
                                "description": "本页要点，可多行，每行一个要点",
                            },
                        },
                        "required": ["title", "content"],
                    },
                },
            },
            "required": ["title", "slides"],
        },
    },
}

PPT_TOOLS = [GET_BUSINESS_DATA_TOOL, GENERATE_PPTX_TOOL]

SYSTEM_MESSAGE = """你是 PPT 汇报助手，具备以下能力：
1. 根据用户需求用自然语言查询业务数据（如本月销售数据、区域/产品统计等），得到结构化结果。
2. 根据查询到的数据或用户给的大纲，组织成汇报用幻灯片（标题 + 要点），并生成 PPT 文件。

当用户说「基于本月销售数据做个汇报PPT」或类似需求时，你应该：
- 先调用 get_business_data 查询相关数据（如：本月销售数据、按产品/区域汇总等）；
- 根据返回的数据总结出汇报结构，再调用 generate_pptx，传入汇报标题和 slides（每页标题与要点）；
- 最后明确告知用户 PPT 已生成，并给出下载/预览链接（file_id 对应的预览地址）。

slides 的 content 建议简洁有力，每条一行，便于观众阅读。若数据中有具体数字，可写入要点中。"""


def _get_business_data(question: str) -> dict:
    """调用 Text2SQL 获取业务数据。"""
    from service.ai.text2sql import text2sql_run

    if not (question or "").strip():
        return {"error": "请提供查询描述", "data": [], "sql": ""}
    out = text2sql_run(question=question.strip(), model="qwen-turbo", max_rows=500)
    if out.get("error"):
        return {
            "error": out["error"],
            "sql": out.get("sql", ""),
            "data": out.get("data", []),
        }
    # 返回摘要 + 数据，便于模型组织 PPT
    data = out.get("data") or []
    summary = f"查询成功，共 {len(data)} 条记录。SQL: {out.get('sql', '')}"
    return {"summary": summary, "sql": out.get("sql", ""), "data": data}


def _generate_pptx_save(title: str, slides: List[Dict[str, Any]]) -> dict:
    """用 python-pptx 生成 PPT 并保存到上传目录，写入 manifest，返回 file_id 与预览地址。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt

    from service.ai.files import (
        _ensure_upload_dirs,
        _load_manifest,
        _save_manifest,
        _upload_root,
    )

    if not title:
        return {"error": "请提供 title"}
    slides = slides or []
    file_id = uuid.uuid4().hex
    root = _ensure_upload_dirs()
    ppt_dir = os.path.join(root, "ppt")
    os.makedirs(ppt_dir, exist_ok=True)
    # 磁盘文件名只用 file_id 与固定后缀，避免中文/特殊字符
    filename = f"{file_id}_report.pptx"
    rel_path = os.path.join("ppt", filename)
    abs_path = os.path.join(root, rel_path)

    prs = Presentation()
    # 封面：标题
    title_layout = prs.slide_layouts[0]
    cover = prs.slides.add_slide(title_layout)
    cover.shapes.title.text = title
    # 内容页
    content_layout = prs.slide_layouts[1]  # title and content
    for s in slides:
        slide_title = (s.get("title") or "").strip() or "（无标题）"
        content = (s.get("content") or "").strip()
        slide = prs.slides.add_slide(content_layout)
        slide.shapes.title.text = slide_title
        tf = slide.placeholders[1].text_frame
        tf.clear()
        for line in content.split("\n"):
            line = line.strip()
            if line:
                p = tf.add_paragraph()
                p.text = line
                p.space_after = Pt(6)
    prs.save(abs_path)
    size = os.path.getsize(abs_path)

    display_name = (title[:80] + "..") if len(title) > 80 else title
    if not display_name.strip():
        display_name = "汇报.pptx"
    elif not display_name.endswith(".pptx"):
        display_name = display_name + ".pptx"
    manifest = _load_manifest()
    manifest[file_id] = {
        "path": rel_path,
        "name": display_name,
        "size": size,
        "kb_id": "ppt",
        "created_at": __import__("datetime").datetime.utcnow().isoformat() + "Z",
    }
    _save_manifest(manifest)

    preview_url = f"/ai/files/{file_id}/preview"
    return {
        "file_id": file_id,
        "preview_url": preview_url,
        "message": f"PPT 已生成，共 {1 + len(slides)} 页。可访问 preview_url 预览或下载。",
    }


def _run_tool(name: str, arguments: dict) -> str:
    """执行工具，返回字符串结果。"""
    if name == "get_business_data":
        q = arguments.get("question", "")
        out = _get_business_data(q)
        return json.dumps(out, ensure_ascii=False, default=str)
    if name == "generate_pptx":
        title = arguments.get("title", "")
        slides = arguments.get("slides", [])
        out = _generate_pptx_save(title=title, slides=slides)
        return json.dumps(out, ensure_ascii=False, default=str)
    return json.dumps({"error": f"未知工具: {name}"}, ensure_ascii=False)


def run_mcp_ppt_chat(
    messages: List[dict],
    model: str = "qwen-turbo",
    system_message: Optional[str] = None,
    max_iterations: int = 8,
) -> dict:
    """
    执行 PPT 助手对话：先查数据再生成 PPT，返回 history、steps、final_answer（含 file_id/预览链接）。
    """
    sys_msg = system_message if system_message is not None else SYSTEM_MESSAGE
    current = [{"role": "system", "content": sys_msg}] + [dict(m) for m in messages]
    reply_messages: List[Dict[str, Any]] = []
    steps: List[Dict[str, Any]] = []
    final_answer = ""

    for _ in range(max_iterations):
        resp = dashscope.Generation.call(
            model=model,
            messages=current,
            tools=PPT_TOOLS,
            tool_choice="auto",
        )
        if resp.status_code != HTTPStatus.OK:
            return {
                "error": f"{getattr(resp, 'code', resp.status_code)} - {getattr(resp, 'message', '')}",
                "reply_messages": reply_messages,
                "steps": steps,
                "history": current,
                "final_answer": final_answer,
            }

        msg = resp.output.choices[0].message
        content = getattr(msg, "content", None) or ""
        tool_calls = getattr(msg, "tool_calls", None) or []

        if not tool_calls:
            final_answer = (content or "").strip()
            if final_answer:
                reply_messages.append({"role": "assistant", "content": final_answer})
                steps.append({"type": "assistant", "content": final_answer})
                current.append({"role": "assistant", "content": final_answer})
            break

        tcs_for_history = []
        function_messages = []
        for tc in tool_calls:
            fn = (
                tc.get("function", {})
                if isinstance(tc, dict)
                else getattr(tc, "function", {})
            )
            if isinstance(fn, dict):
                fn_name = fn.get("name", "")
                fn_args_str = fn.get("arguments", "{}")
            else:
                fn_name = getattr(fn, "name", "")
                fn_args_str = getattr(fn, "arguments", "{}")
            try:
                fn_args = (
                    json.loads(fn_args_str)
                    if isinstance(fn_args_str, str)
                    else fn_args_str
                )
            except Exception:
                fn_args = {}

            steps.append(
                {
                    "type": "tool_call_start",
                    "tool_name": fn_name,
                    "arguments": fn_args_str,
                }
            )
            result = _run_tool(fn_name, fn_args)
            steps.append(
                {"type": "tool_call_end", "tool_name": fn_name, "output": result}
            )

            reply_messages.append(
                {
                    "role": "assistant",
                    "content": f'▼ 调用工具 "{fn_name}"\n{fn_args_str}',
                }
            )
            reply_messages.append(
                {
                    "role": "function",
                    "name": fn_name,
                    "content": f"▼ 工具返回\n{result}",
                }
            )

            tcs_for_history.append(
                tc
                if isinstance(tc, dict)
                else {
                    "id": getattr(tc, "id", ""),
                    "type": "function",
                    "function": {"name": fn_name, "arguments": fn_args_str},
                }
            )
            function_messages.append(
                {"role": "function", "name": fn_name, "content": result}
            )

        current.append(
            {"role": "assistant", "content": content or "", "tool_calls": tcs_for_history}
        )
        current.extend(function_messages)

    return {
        "reply_messages": reply_messages,
        "steps": steps,
        "history": current,
        "final_answer": final_answer,
    }


def get_mcp_ppt_info() -> dict:
    """返回 PPT 助手元信息与工具列表，供低代码/云平台展示与接入。"""
    tools_info = []
    for t in PPT_TOOLS:
        fn = t.get("function", {})
        tools_info.append(
            {
                "name": fn.get("name", ""),
                "description": fn.get("description", ""),
                "parameters": fn.get("parameters", {}),
            }
        )
    return {
        "name": "PPT 汇报助手",
        "description": "根据自然语言拉取业务数据并自动生成汇报 PPT，支持「基于本月销售数据做个汇报PPT」等一句话生成。可接入活字格、阿里云百炼等低代码/云平台。",
        "tools": tools_info,
    }
